from pptx import Presentation
import csv
import unittest
import codecs
import os
import streamlit as st
import numpy as np
import pandas as pd
url ="https://www.linkedin.com/in/alioune-gaye-1a5161172/"
def main():
    st.title("Fonds d’Etude et d’Aide au Secteur Privé (FASEP)")
    st.image("ministre.PNG")
    st.sidebar.write("<p style='text-align: center;'> Alioune Gaye:  %s</p>" % url, unsafe_allow_html=True)
    st.sidebar.write("Application qui permet de récupérer et de stocker dans un fichier csv les dates de signature de la convention de don de la subvention du fonds d’étude et d’aide au secteur privé (FASEP), le montant de la subvention du FASEP, et l’avis du service économique de l’ambassade pour le premier terme intermédiaire de la subvention à partir d'un document powerpoint « Exemple FASEP.pptx » .")
    def extraction_de_data(pptx_file, fichier_sortie):
        data=[]
        try:
            power = Presentation(pptx_file)
            # On choisit un fichier en mode écriture au format CSV pour écrire les données extraites. 
            with open(fichier_sortie, "w", newline='', encoding="utf-8-sig") as fichier_csv:
                # Initialisation d'un lecteur pour écrire dans le fichier CSV.
                csv_file = csv.writer(fichier_csv)
                # On va parcourir toutes les diapositives du fichier pptx
                for slide in power.slides:
                    # Toutes les formes présentes sur chaque diapositive. 
                    for shape in slide.shapes:
                        # Comme les informations qui nous intéressent sont sur des tableaux. Alors, on vérifie si la forme est un tableau.
                        if shape.has_table:
                            # Maintenant, on va accéder à chaque tableau et parcourir chaque ligne du tableau.
                            tableau = shape.table
                            for row in tableau.rows:
                                trouve = False
                                # On va parcourir chaque cellule de la ligne, puis récupérer le texte.
                                for cell in row.cells:
                                    texte_cellule = cell.text
                                    # On dois vérifier si l'une des phrases recherchées est présente dans le texte de la cellule
                                    if "Date de signature de la convention Natixis" in texte_cellule or \
                                       "Montant et date de paiement de l'acompte" in texte_cellule or \
                                       "Date de l'avis" in texte_cellule or \
                                       "Avis sur le versement intermédiaire" in texte_cellule:
                                        trouve = True
                                        break 
                                # Si l'une des phrases est trouvée dans la ligne, on récupére la ligne entière
                                if trouve:
                                    # On récupère le texte de chaque cellule dans la ligne et les ajoute à une liste. Puis, on ajoute la ligne du tableau dans le fichier CSV
                                    ligne_tableau = [cell.text if cell.text else "" for cell in row.cells]
                                    csv_file.writerow(ligne_tableau)
                                    data.append(ligne_tableau)
            st.write("Extraction des informations terminée avec succès.")                        
            #Un dataFrame à partir des données extraites
            df = pd.DataFrame(data)                        
            st.dataframe(df.head(4))                                  
        except Exception as e:
            st.write(f"Une erreur s'est produite, recommencez : votre fichier n'est pas conforme ou n'est pas de type pptx ou ppt: {str(e)}")    
    uploaded_file = st.file_uploader("Importez le fichier Powerpoint", type=['ppt', 'pptx'])
    if uploaded_file is not None:
        extraction_de_data(uploaded_file, "fichier_sortie.csv")
        with open("fichier_sortie.csv",'rb') as f :
            st.download_button("Cliquez ici pour télécharger les informations recueillies." , f, file_name="fichier_sortie.csv")
    else:
        st.write("Veuillez télécharger un fichier de type PowerPoint svp.")
if __name__ == '__main__':
    main()